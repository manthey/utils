#!/usr/bin/env python3
# /// script
# requires-python = '>=3.12'
# dependencies = [
#     'pillow',
#     'python-pptx',
#     'openai',
# ]
# ///
import argparse
import base64
import io
import logging
import os
import sys
from pathlib import Path
from typing import Any

import PIL.Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
logger.addHandler(logging.NullHandler())


def describe_image(
    url: str, model: str, b64_image: str, system: str, user: str,
    options: dict[str, Any] | None = None, is_png: bool = False,
) -> str:
    import openai

    client = openai.OpenAI(base_url=f'{url}/v1', api_key='ollama', timeout=300)
    messages = [{
        'role': 'system',
        'content': [{'type': 'text', 'text': system}],
    }, {
        'role': 'user',
        'content': [{
            'type': 'text',
            'text': user,
        }, {
            'type': 'image_url',
            'image_url': {'url': f'data:image/jpeg;base64,{b64_image}'},
        }],
    }]
    if not system:
        messages[0:1] = []
    response = client.chat.completions.create(
        model=model, messages=messages, **(options or {}))
    message = response.choices[0].message.content
    if '```' in message:
        message = message.split('```')[1].split('\n', 1)[-1]
    return message


def extract_text(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    paragraphs = [p for p in paragraphs if p != '‹#›']
    return ['\n'.join(paragraphs)] if paragraphs else []


def extract_table(shape) -> list[str]:
    if not shape.has_table:
        return []
    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(' | '.join(cells))
    return ['\n'.join(rows)] if rows else []


def process_slide(slide, slide_index: int, args) -> str:
    lines = [f'## Slide {slide_index}\n']
    text_blocks = []
    image_index = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_blocks.extend(extract_text(shape))
        if shape.has_table:
            text_blocks.extend(extract_table(shape))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_index += 1
            img = shape.image.blob
            is_png = img.startswith(b'\x89PNG')
            if not is_png and img[:1] != b'\xff':
                img = PIL.Image.open(io.BytesIO(img))
                img = img.convert('L' if img.mode in {'L', 'LA'} else 'RGB')
                buf = io.BytesIO()
                img.save(buf, format='JPEG', quality=95)
                img = buf.getvalue()
            b64 = base64.b64encode(img).decode('utf-8')
            description = describe_image(
                url=args.url,
                model=args.model,
                b64_image=b64,
                system=args.system,
                user=args.user,
                is_png=is_png,
            )
            text_blocks.append(f'**Image {image_index}:** {description}')
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.has_text_frame:
                    text_blocks.extend(extract_text(child))
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_index += 1
                    b64 = base64.b64encode(child.image.blob).decode('utf-8')
                    description = describe_image(
                        url=args.url,
                        model=args.model,
                        b64_image=b64,
                        system=args.system,
                        user=args.user,
                    )
                    text_blocks.append(f'**Image {image_index}:** {description}')
    if text_blocks:
        lines.append('\n\n'.join(text_blocks))
    return '\n'.join(lines) + '\n'


def process_file(filepath, args):
    filepath = Path(filepath)
    presentation = Presentation(filepath)
    desc = [f'# {filepath.name}\n']
    for index, slide in enumerate(presentation.slides, start=1):
        desc += [process_slide(slide, index, args)]
    return '\n'.join(desc)


def process_directory(args):  # noqa
    suffix = f'.{args.suffix.lstrip(".")}'
    for input_path in args.inputs:
        target = Path(input_path)
        if target.is_file():
            file_list = [target]
        elif target.is_dir():
            file_list = sorted(target.rglob('*')) if args.recurse else sorted(target.iterdir())
        else:
            continue
        for filepath in file_list:
            if not filepath.is_file():
                continue
            if not str(filepath).endswith(('.pptx', '.ppt')) and filepath not in args.inputs:
                continue
            md_path = filepath.with_suffix(suffix)
            if args.out:
                if os.path.isdir(args.out):
                    md_path = Path(args.out) / md_path.name
                else:
                    md_path = Path(args.out)
            if (not args.overwrite and md_path.exists() and
                    md_path.stat().st_mtime > filepath.stat().st_mtime):
                continue
            if args.out and not os.path.isdir(args.out):
                args.overwrite = False
            if args.list:
                print(f'{filepath} -> {md_path}')
                continue
            try:
                print(filepath)
                description = process_file(filepath, args)
                print(description)
                if not args.dry_run:
                    md_path.write_text(description, encoding='utf-8')
                    print(f'Created {md_path.name}')
                else:
                    print(f'Would have created {md_path.name}')
            except Exception as exc:
                print(f'Failed processing {filepath.name}: {exc}')
                if args.raise_errors:
                    raise


def main():
    parser = argparse.ArgumentParser(
        description='Convert PPTX to markdown with image descriptions.')
    parser.add_argument(
        'inputs', nargs='+',
        help='One or more files or directories to process.')
    parser.add_argument(
        '--recurse', '-r', action='store_true',
        help='Recurse into input directories')
    parser.add_argument(
        '--suffix', '--ext', default='.description.md',
        help='File extension to use for description files.')
    parser.add_argument(
        '--out', '--output',
        help='If an existing directory, the location to store outputs.  If a '
        'single path or non-existent path, write the first description to '
        'this file and then stop.')
    parser.add_argument(
        '--url', default=os.environ.get('OLLAMA_BASE_URL', 'http://localhost:11434'),
        help='Ollama base URL.  Default %(default)s.')
    parser.add_argument(
        '--api-key', default='ollama',
        help='API key sent to the endpoint.  Default %(default)s.')
    parser.add_argument(
        '--model', '-m', default='qwen2.5vl:7b',
        help='Vision model identifier.  Default %(default)s.')
    parser.add_argument(
        '--system',
        default='You describe images concisely for document summarization '
        'and search retrieval.  You never use emojis, slang, or metaphors.',
        help='System prompt for image description')
    parser.add_argument(
        '--user', default='Describe this image in detail.',
        help='User prompt for image description')
    parser.add_argument(
        '--overwrite', '-y', action='store_true',
        help='Overwrite existing companion markdown files')
    parser.add_argument(
        '-n', '--dry-run', action='store_true',
        help='Do not actually write markdown files')
    parser.add_argument(
        '--list', '-l', action='store_true',
        help='Just list what files would be processed without actually doing anything.')
    parser.add_argument(
        '--raise', dest='raise_errors', action='store_true',
        help='Raise on errors instead of ignoring them.')
    parser.add_argument(
        '--verbose', '-v', action='count', default=0,
        help='Increase verbosity')
    args = parser.parse_args()
    logger.setLevel(max(1, logging.WARNING - args.verbose * 10))
    logger.addHandler(logging.StreamHandler(sys.stderr))
    logger.debug('Parsed arguments: %r', args)
    process_directory(args)


if __name__ == '__main__':
    main()
